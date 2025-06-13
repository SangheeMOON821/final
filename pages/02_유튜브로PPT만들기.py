import streamlit as st
from youtube_transcript_api import YouTubeTranscriptApi, TranscriptsDisabled
from youtube_transcript_api._errors import RequestBlocked
from pptx import Presentation
from pptx.util import Inches
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
import openai
import re
import tempfile
import os
from dotenv import load_dotenv  # ✅ dotenv 추가

# ✅ .env 파일 로드
load_dotenv()

# ✅ 환경변수에서 API 키 불러오기
YOUTUBE_API_KEY = os.getenv("AIzaSyCiOYEp0EsDCRl2xn5exfxpJyv78SJfIOQ")
openai.api_key = os.getenv("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")

# --------------------------------------
def extract_video_id(url):
    regex = r"(?:v=|\/)([0-9A-Za-z_-]{11}).*"
    match = re.search(regex, url)
    return match.group(1) if match else None

def get_transcript(video_id):
    try:
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        return " ".join([item["text"] for item in transcript])
    except TranscriptsDisabled:
        return None
    except RequestBlocked:
        return "BLOCKED"
    except Exception:
        return None

def get_video_metadata(video_id):
    try:
        youtube = build("youtube", "v3", developerKey=YOUTUBE_API_KEY)
        request = youtube.videos().list(part="snippet", id=video_id)
        response = request.execute()

        if not response["items"]:
            return None, None
        item = response["items"][0]["snippet"]
        return item["title"], item.get("description", "")
    except HttpError as e:
        st.error(f"[YouTube API Error] {e}")
        return None, None

def summarize_text(text, max_chars=800):
    sentences = text.split('. ')
    chunks = []
    current_chunk = ""
    for sentence in sentences:
        if len(current_chunk) + len(sentence) < max_chars:
            current_chunk += sentence + ". "
        else:
            chunks.append(current_chunk.strip())
            current_chunk = sentence + ". "
    if current_chunk:
        chunks.append(current_chunk.strip())
    return chunks

def summarize_with_gpt(text, lang="ko"):
    prompt = f"다음 유튜브 영상 소개 내용을 기반으로 주요 내용을 {lang}로 슬라이드처럼 요약해줘:\n\n{text}\n\n요약:"
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=600,
            temperature=0.7
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"[GPT Error] {e}")
        return "GPT 요약에 실패했습니다."

def create_ppt(youtube_url, title, summaries):
    prs = Presentation()

    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title or "YouTube Summary Presentation"
    slide.placeholders[1].text = f"Link: {youtube_url}"

    content_layout = prs.slide_layouts[1]
    for idx, summary in enumerate(summaries, start=1):
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = f"요약 {idx}"
        slide.placeholders[1].text = summary

    tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp_file.name)
    return tmp_file.name

# --------------------------------------
st.set_page_config(page_title="유튜브 요약 PPT 생성기")
st.title("🎞️ 유튜브 영상 요약 PPT 생성기")

youtube_url = st.text_input("유튜브 링크를 입력하세요")

if st.button("📄 PPT 만들기"):
    if not youtube_url:
        st.warning("유튜브 링크를 입력해주세요.")
        st.stop()

    video_id = extract_video_id(youtube_url)
    if not video_id:
        st.error("유효한 유튜브 링크가 아닙니다.")
        st.stop()

    with st.spinner("1단계: 유튜브 영상 정보 확인 중..."):
        title, description = get_video_metadata(video_id)
        if title is None:
            st.error("YouTube API 요청 실패 - API 키 확인 또는 영상 권한 확인 필요.")
            st.stop()

    with st.spinner("2단계: 자막 정보 가져오는 중..."):
        transcript = get_transcript(video_id)

    if transcript == "BLOCKED" or transcript is None:
        st.warning("자막이 없거나 차단되어 GPT로 요약합니다.")
        if not description:
            st.error("영상 설명이 없어 요약이 불가능합니다.")
            st.stop()
        with st.spinner("GPT 요약 생성 중..."):
            summary_text = summarize_with_gpt(description)
            summaries = summary_text.split("\n")
    else:
        with st.spinner("자막 기반 요약 생성 중..."):
            summaries = summarize_text(transcript)

    with st.spinner("3단계: PPT 생성 중..."):
        ppt_path = create_ppt(youtube_url, title, summaries)

    with open(ppt_path, "rb") as f:
        st.success("🎉 PPT 생성 완료!")
        st.download_button(
            label="📥 PPT 다운로드",
            data=f,
            file_name="youtube_summary.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    os.remove(ppt_path)
