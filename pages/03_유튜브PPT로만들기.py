import streamlit as st
from pytube import YouTube
from youtube_transcript_api import YouTubeTranscriptApi, TranscriptsDisabled, NoTranscriptAvailable
from pptx import Presentation
from pptx.util import Inches
from urllib.parse import urlparse, parse_qs
import textwrap
import os
import whisper

def extract_video_id(url):
    parsed_url = urlparse(url)
    if parsed_url.hostname in ["www.youtube.com", "youtube.com"]:
        return parse_qs(parsed_url.query).get("v", [None])[0]
    elif parsed_url.hostname == "youtu.be":
        return parsed_url.path[1:]
    return None

def get_youtube_transcript(video_id):
    try:
        # 가능한 자막 언어 리스트 가져오기
        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
        transcript = transcript_list.find_transcript(['ko', 'en'])
        return " ".join([t["text"] for t in transcript.fetch()])
    except (TranscriptsDisabled, NoTranscriptAvailable):
        return None
    except Exception as e:
        return None

def download_audio_from_youtube(video_url, output_path="audio.mp4"):
    yt = YouTube(video_url)
    stream = yt.streams.filter(only_audio=True).first()
    stream.download(filename=output_path)
    return output_path

def transcribe_audio_whisper(audio_path):
    model = whisper.load_model("base")  # 필요한 경우 'small', 'medium', 'large' 등도 가능
    result = model.transcribe(audio_path, language='ko')
    return result["text"]

def summarize_text(text, max_sentences=5):
    sentences = text.split(". ")
    return ". ".join(sentences[:max_sentences]) + ("." if sentences else "")

def split_summary(summary, max_length=500):
    return textwrap.wrap(summary, width=max_length)

def create_ppt(youtube_url, summary_chunks):
    prs = Presentation()

    # 첫 슬라이드: 링크 정보
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "YouTube 영상 요약"
    subtitle.text = youtube_url

    # 요약 슬라이드들
    for i, chunk in enumerate(summary_chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = f"요약 {i+1}"
        content.text = chunk

    return prs

st.title("📺 유튜브 영상으로 PPT 요약 만들기")

youtube_url = st.text_input("유튜브 링크를 입력하세요:")

if youtube_url:
    video_id = extract_video_id(youtube_url)
    if not video_id:
        st.error("유효한 유튜브 링크가 아닙니다.")
    else:
        st.info("자막 확인 중입니다...")
        transcript_text = get_youtube_transcript(video_id)

        if transcript_text is None:
            st.warning("자막이 없습니다. Whisper로 음성을 인식 중입니다...")
            audio_path = download_audio_from_youtube(youtube_url)
            transcript_text = transcribe_audio_whisper(audio_path)
            os.remove(audio_path)
            st.success("Whisper로 자막 생성 완료!")

        summary = summarize_text(transcript_text, max_sentences=15)
        summary_chunks = split_summary(summary)

        prs = create_ppt(youtube_url, summary_chunks)

        output_path = "youtube_summary.pptx"
        prs.save(output_path)

        with open(output_path, "rb") as f:
            st.download_button(
                label="📥 PPT 다운로드",
                data=f,
                file_name="youtube_summary.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        os.remove(output_path)
